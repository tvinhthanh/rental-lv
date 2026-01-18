"""
Script để tạo PowerPoint báo cáo đề tài Hệ thống Quản lý Cho thuê Xe
Cần cài đặt: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    """Tạo presentation mới"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Trang bìa
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 23, 42)  # slate-900
    
    # Title
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "HỆ THỐNG QUẢN LÝ CHO THUÊ XE"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Car Rental Management System"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RGBColor(148, 163, 184)  # slate-400
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Author info
    author_box = slide1.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(2))
    author_frame = author_box.text_frame
    author_frame.text = "Sinh viên: [Tên sinh viên]\nGiảng viên hướng dẫn: [Tên giảng viên]\nNgày: [Ngày báo cáo]"
    author_para = author_frame.paragraphs[0]
    author_para.font.size = Pt(18)
    author_para.font.color.rgb = RGBColor(203, 213, 225)  # slate-300
    author_para.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Mục lục
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    title2 = slide2.shapes.title
    title2.text = "MỤC LỤC"
    content2 = slide2.placeholders[1]
    tf2 = content2.text_frame
    tf2.text = "1. Giới thiệu đề tài"
    p2 = tf2.add_paragraph()
    p2.text = "2. Mục tiêu và phạm vi"
    p2 = tf2.add_paragraph()
    p2.text = "3. Phân tích yêu cầu"
    p2 = tf2.add_paragraph()
    p2.text = "4. Kiến trúc hệ thống"
    p2 = tf2.add_paragraph()
    p2.text = "5. Công nghệ sử dụng"
    p2 = tf2.add_paragraph()
    p2.text = "6. Cơ sở dữ liệu"
    p2 = tf2.add_paragraph()
    p2.text = "7. Chức năng chính"
    p2 = tf2.add_paragraph()
    p2.text = "8. Giao diện người dùng"
    p2 = tf2.add_paragraph()
    p2.text = "9. Kết quả đạt được"
    p2 = tf2.add_paragraph()
    p2.text = "10. Kết luận và hướng phát triển"
    
    # Slide 3: Giới thiệu đề tài
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    title3.text = "GIỚI THIỆU ĐỀ TÀI"
    content3 = slide3.placeholders[1]
    tf3 = content3.text_frame
    tf3.text = "Vấn đề thực tế:"
    p3 = tf3.add_paragraph()
    p3.text = "• Nhu cầu thuê xe ngày càng tăng"
    p3 = tf3.add_paragraph()
    p3.text = "• Quản lý thủ công gặp nhiều khó khăn"
    p3 = tf3.add_paragraph()
    p3.text = "• Cần hệ thống tự động hóa quy trình"
    p3 = tf3.add_paragraph()
    p3.text = ""
    p3 = tf3.add_paragraph()
    p3.text = "Giải pháp:"
    p3 = tf3.add_paragraph()
    p3.text = "• Xây dựng hệ thống quản lý cho thuê xe toàn diện"
    p3 = tf3.add_paragraph()
    p3.text = "• Tự động hóa quy trình từ đặt xe đến thanh toán"
    p3 = tf3.add_paragraph()
    p3.text = "• Hỗ trợ nhiều vai trò: Admin, Nhân viên, Khách hàng"
    
    # Slide 4: Mục tiêu và phạm vi
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title4 = slide4.shapes.title
    title4.text = "MỤC TIÊU VÀ PHẠM VI"
    content4 = slide4.placeholders[1]
    tf4 = content4.text_frame
    tf4.text = "Mục tiêu:"
    p4 = tf4.add_paragraph()
    p4.text = "✅ Xây dựng hệ thống quản lý cho thuê xe đầy đủ"
    p4 = tf4.add_paragraph()
    p4.text = "✅ Tự động hóa quy trình booking → contract → handover → return → invoice"
    p4 = tf4.add_paragraph()
    p4.text = "✅ Hỗ trợ thanh toán online (Stripe) và tiền mặt"
    p4 = tf4.add_paragraph()
    p4.text = "✅ Quản lý đa chi nhánh"
    p4 = tf4.add_paragraph()
    p4.text = "✅ SEO và Content Marketing (Blog, Pages)"
    p4 = tf4.add_paragraph()
    p4.text = "✅ CRM và Marketing (Promotions, Loyalty Program)"
    p4 = tf4.add_paragraph()
    p4.text = ""
    p4 = tf4.add_paragraph()
    p4.text = "Phạm vi:"
    p4 = tf4.add_paragraph()
    p4.text = "• Backend: NestJS + TypeScript + MongoDB"
    p4 = tf4.add_paragraph()
    p4.text = "• Frontend: Next.js 14 + React 18 + Tailwind CSS"
    p4 = tf4.add_paragraph()
    p4.text = "• Payment: Stripe + Cash Payment"
    p4 = tf4.add_paragraph()
    p4.text = "• Real-time: Socket.io cho notifications"
    
    # Slide 5: Phân tích yêu cầu
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title5 = slide5.shapes.title
    title5.text = "PHÂN TÍCH YÊU CẦU"
    content5 = slide5.placeholders[1]
    tf5 = content5.text_frame
    tf5.text = "Yêu cầu chức năng:"
    p5 = tf5.add_paragraph()
    p5.text = "1. Quản lý xe: CRUD xe, danh mục, thương hiệu, giá"
    p5 = tf5.add_paragraph()
    p5.text = "2. Quản lý đặt xe: Booking, kiểm tra lịch trống, tính giá tự động"
    p5 = tf5.add_paragraph()
    p5.text = "3. Quản lý hợp đồng: Tạo hợp đồng, bàn giao, nhận xe"
    p5 = tf5.add_paragraph()
    p5.text = "4. Quản lý thanh toán: Invoice, Payment (Stripe/Cash), Deposit"
    p5 = tf5.add_paragraph()
    p5.text = "5. Quản lý khách hàng: Customer, Membership tiers, Loyalty points"
    p5 = tf5.add_paragraph()
    p5.text = "6. Quản lý nhân viên: Employee, Branch, Permissions"
    p5 = tf5.add_paragraph()
    p5.text = "7. Marketing: Promotions, Blog, SEO pages"
    p5 = tf5.add_paragraph()
    p5.text = "8. Báo cáo: Dashboard, Statistics, Audit logs"
    p5 = tf5.add_paragraph()
    p5.text = ""
    p5 = tf5.add_paragraph()
    p5.text = "Yêu cầu phi chức năng:"
    p5 = tf5.add_paragraph()
    p5.text = "• Bảo mật: JWT Authentication, Role-based Access Control"
    p5 = tf5.add_paragraph()
    p5.text = "• Hiệu năng: Caching với Redis, Optimized queries"
    p5 = tf5.add_paragraph()
    p5.text = "• Scalability: Microservices-ready architecture"
    p5 = tf5.add_paragraph()
    p5.text = "• Responsive: Mobile-friendly UI"
    
    # Slide 6: Kiến trúc hệ thống
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    title6 = slide6.shapes.title
    title6.text = "KIẾN TRÚC HỆ THỐNG"
    content6 = slide6.placeholders[1]
    tf6 = content6.text_frame
    tf6.text = "Kiến trúc 3-Layer:"
    p6 = tf6.add_paragraph()
    p6.text = ""
    p6 = tf6.add_paragraph()
    p6.text = "1. Presentation Layer: Next.js Frontend (Admin/Employee/User portals)"
    p6 = tf6.add_paragraph()
    p6.text = "2. Business Logic Layer: NestJS Backend (Modules, Services, Controllers)"
    p6 = tf6.add_paragraph()
    p6.text = "3. Data Layer: MongoDB + Prisma ORM"
    p6 = tf6.add_paragraph()
    p6.text = ""
    p6 = tf6.add_paragraph()
    p6.text = "Communication:"
    p6 = tf6.add_paragraph()
    p6.text = "• HTTP/REST API giữa Frontend và Backend"
    p6 = tf6.add_paragraph()
    p6.text = "• WebSocket (Socket.io) cho real-time notifications"
    p6 = tf6.add_paragraph()
    p6.text = "• Prisma Client cho database access"
    
    # Slide 7: Công nghệ Backend
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    title7 = slide7.shapes.title
    title7.text = "CÔNG NGHỆ SỬ DỤNG - BACKEND"
    content7 = slide7.placeholders[1]
    tf7 = content7.text_frame
    tf7.text = "Core Framework:"
    p7 = tf7.add_paragraph()
    p7.text = "• NestJS 10.x - Progressive Node.js framework"
    p7 = tf7.add_paragraph()
    p7.text = "• TypeScript - Type-safe development"
    p7 = tf7.add_paragraph()
    p7.text = "• Prisma - Modern ORM for MongoDB"
    p7 = tf7.add_paragraph()
    p7.text = ""
    p7 = tf7.add_paragraph()
    p7.text = "Authentication & Security:"
    p7 = tf7.add_paragraph()
    p7.text = "• JWT - JSON Web Tokens"
    p7 = tf7.add_paragraph()
    p7.text = "• Passport - Authentication middleware"
    p7 = tf7.add_paragraph()
    p7.text = "• bcryptjs - Password hashing"
    p7 = tf7.add_paragraph()
    p7.text = "• class-validator - DTO validation"
    p7 = tf7.add_paragraph()
    p7.text = ""
    p7 = tf7.add_paragraph()
    p7.text = "Additional Services:"
    p7 = tf7.add_paragraph()
    p7.text = "• Socket.io - Real-time notifications"
    p7 = tf7.add_paragraph()
    p7.text = "• Cloudinary - Image upload & storage"
    p7 = tf7.add_paragraph()
    p7.text = "• Stripe - Payment gateway"
    p7 = tf7.add_paragraph()
    p7.text = "• Redis - Caching layer"
    
    # Slide 8: Công nghệ Frontend
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    title8 = slide8.shapes.title
    title8.text = "CÔNG NGHỆ SỬ DỤNG - FRONTEND"
    content8 = slide8.placeholders[1]
    tf8 = content8.text_frame
    tf8.text = "Core Framework:"
    p8 = tf8.add_paragraph()
    p8.text = "• Next.js 14 - React framework với App Router"
    p8 = tf8.add_paragraph()
    p8.text = "• React 18 - UI library"
    p8 = tf8.add_paragraph()
    p8.text = "• TypeScript - Type safety"
    p8 = tf8.add_paragraph()
    p8.text = ""
    p8 = tf8.add_paragraph()
    p8.text = "UI & Styling:"
    p8 = tf8.add_paragraph()
    p8.text = "• Tailwind CSS - Utility-first CSS framework"
    p8 = tf8.add_paragraph()
    p8.text = "• shadcn/ui - Component library"
    p8 = tf8.add_paragraph()
    p8.text = "• Lucide React - Icon library"
    p8 = tf8.add_paragraph()
    p8.text = ""
    p8 = tf8.add_paragraph()
    p8.text = "State Management:"
    p8 = tf8.add_paragraph()
    p8.text = "• React Query - Server state management"
    p8 = tf8.add_paragraph()
    p8.text = "• Redux Toolkit - Client state management"
    p8 = tf8.add_paragraph()
    p8.text = ""
    p8 = tf8.add_paragraph()
    p8.text = "Forms & Validation:"
    p8 = tf8.add_paragraph()
    p8.text = "• React Hook Form - Form management"
    p8 = tf8.add_paragraph()
    p8.text = "• Zod - Schema validation"
    
    # Slide 9: Cơ sở dữ liệu Overview
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    title9 = slide9.shapes.title
    title9.text = "CƠ SỞ DỮ LIỆU - OVERVIEW"
    content9 = slide9.placeholders[1]
    tf9 = content9.text_frame
    tf9.text = "Database: MongoDB (NoSQL Document Database)"
    p9 = tf9.add_paragraph()
    p9.text = ""
    p9 = tf9.add_paragraph()
    p9.text = "37 Tables/Collections (100% complete)"
    p9 = tf9.add_paragraph()
    p9.text = ""
    p9 = tf9.add_paragraph()
    p9.text = "Phân loại Collections:"
    p9 = tf9.add_paragraph()
    p9.text = "🔴 Core Business (20): User, Customer, Employee, Vehicle, Booking, Contract, Invoice, Payment..."
    p9 = tf9.add_paragraph()
    p9.text = "🟡 SEO & Content (5): Blog, Page, SeoRedirect, VehicleCategory, Brand"
    p9 = tf9.add_paragraph()
    p9.text = "🟢 Marketing & CRM (6): Promotion, LoyaltyProgram, MarketingCampaign, CustomerSegment..."
    p9 = tf9.add_paragraph()
    p9.text = "🟣 Enterprise (6): Tenant, SubscriptionPlan, Partner, AuditLog, Settings..."
    p9 = tf9.add_paragraph()
    p9 = tf9.add_paragraph()
    p9.text = "Features:"
    p9 = tf9.add_paragraph()
    p9.text = "• Prisma Schema - Type-safe database access"
    p9 = tf9.add_paragraph()
    p9.text = "• Referential Integrity với Prisma relations"
    p9 = tf9.add_paragraph()
    p9.text = "• Indexes cho performance"
    p9 = tf9.add_paragraph()
    p9.text = "• Audit Fields (createdAt, updatedAt)"
    
    # Slide 10: Chức năng chính - Quản lý xe
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    title10 = slide10.shapes.title
    title10.text = "CHỨC NĂNG CHÍNH - QUẢN LÝ XE"
    content10 = slide10.placeholders[1]
    tf10 = content10.text_frame
    tf10.text = "Vehicle Management:"
    p10 = tf10.add_paragraph()
    p10.text = "✅ CRUD xe (thêm, sửa, xóa, xem)"
    p10 = tf10.add_paragraph()
    p10.text = "✅ Upload ảnh (Cloudinary)"
    p10 = tf10.add_paragraph()
    p10.text = "✅ Quản lý danh mục và thương hiệu"
    p10 = tf10.add_paragraph()
    p10.text = "✅ Bảng giá theo loại xe"
    p10 = tf10.add_paragraph()
    p10.text = "✅ Quản lý chi nhánh"
    p10 = tf10.add_paragraph()
    p10.text = "✅ SEO fields (slug, metaTitle, metaDescription)"
    p10 = tf10.add_paragraph()
    p10.text = "✅ ViewCount tracking"
    p10 = tf10.add_paragraph()
    p10.text = ""
    p10 = tf10.add_paragraph()
    p10.text = "Features:"
    p10 = tf10.add_paragraph()
    p10.text = "• Tìm kiếm và lọc xe"
    p10 = tf10.add_paragraph()
    p10.text = "• Kiểm tra lịch trống"
    p10 = tf10.add_paragraph()
    p10.text = "• Quản lý trạng thái (Available, Rented, Maintenance)"
    
    # Thêm các slides còn lại...
    # (Có thể tiếp tục thêm các slides khác tương tự)
    
    # Slide cuối: Cảm ơn
    slide_final = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    background_final = slide_final.background
    fill_final = background_final.fill
    fill_final.solid()
    fill_final.fore_color.rgb = RGBColor(15, 23, 42)
    
    thanks_box = slide_final.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    thanks_frame = thanks_box.text_frame
    thanks_frame.text = "CẢM ƠN!"
    thanks_para = thanks_frame.paragraphs[0]
    thanks_para.font.size = Pt(48)
    thanks_para.font.bold = True
    thanks_para.font.color.rgb = RGBColor(255, 255, 255)
    thanks_para.alignment = PP_ALIGN.CENTER
    
    qa_box = slide_final.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1))
    qa_frame = qa_box.text_frame
    qa_frame.text = "Questions & Answers"
    qa_para = qa_frame.paragraphs[0]
    qa_para.font.size = Pt(24)
    qa_para.font.color.rgb = RGBColor(148, 163, 184)
    qa_para.alignment = PP_ALIGN.CENTER
    
    return prs

if __name__ == "__main__":
    print("Đang tạo PowerPoint báo cáo...")
    presentation = create_presentation()
    output_file = "BaoCao_HeThongQuanLyChoThueXe.pptx"
    presentation.save(output_file)
    print(f"✅ Đã tạo file: {output_file}")
    print(f"📊 Tổng số slides: {len(presentation.slides)}")
